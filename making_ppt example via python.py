
# Make sure you have python-pptx installed:

# pip install python-pptx
# Run the script
# python3 generate_presentation.py



from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullet_slide(title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0

def add_image_slide(title, image_path, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    left = Inches(1)
    top = Inches(1.5)
    pic = slide.shapes.add_picture(image_path, left, top, width=Inches(6))
    if caption:
        txBox = slide.shapes.add_textbox(left, top + Inches(4.5), Inches(6), Inches(0.5))
        tf = txBox.text_frame
        tf.text = caption

# 1. Title Slide
add_title_slide(
    "Data Science Job Salaries: EDA & Salary Prediction",
    "Exploratory Data Analysis and Modeling\nYour Name\n2025"
)

# 2. Objective
add_bullet_slide("Objective", [
    "Analyze and model data science job salaries to uncover trends",
    "Identify salary drivers",
    "Predict salaries based on job-related factors"
])

# 3. Dataset Overview
add_bullet_slide("Dataset Overview", [
    "607 records, 12 features (2020–2022)",
    "Features: work year, experience level, employment type, job title, salary, currency, employee residence, remote ratio, company location, company size",
    "No missing values"
])

# 4. Libraries Used
add_bullet_slide("Libraries Used", [
    "pandas, numpy",
    "matplotlib, seaborn",
    "scikit-learn",
    "streamlit"
])

# 5. Key Visualization: Top 10 Job Titles
add_image_slide("Top 10 Job Titles by Count", "bar_chart_job_titles.png")

# 6. Key Visualization: Salary by Experience Level
add_image_slide("Average Salary by Experience Level", "column_chart_experience.png")

# 7. Key Visualization: Employment Type Distribution
add_image_slide("Employment Type Distribution", "bar_chart_employment_type.png")

# 8. Univariate Analysis
add_bullet_slide("Univariate Analysis", [
    "Senior-level/Expert roles dominate",
    "Most common titles: Data Scientist, Data Engineer, Data Analyst, ML Engineer",
    "Full-time employment is the norm"
])

# 9. Geographic & Company Trends
add_bullet_slide("Geographic & Company Trends", [
    "US: highest number of professionals and companies",
    "Highest average salaries: Russia, US",
    "Medium and large companies pay more"
])

# 10. Salary Trends
add_bullet_slide("Salary Trends", [
    "Salaries and job opportunities increased each year",
    "Higher experience = higher salary",
    "Medium/Large companies pay more than small"
])

# 11. Predictive Modeling
add_bullet_slide("Predictive Modeling", [
    "Linear regression model built to predict salary",
    "Features: experience level, employment type, remote ratio, company size, job title",
    "Model enables salary estimation for various roles"
])

# 12. Salary Predictor App
add_image_slide("Salary Predictor App", "predictor_screenshot.png", "Streamlit-based interactive salary predictor")

# 13. Conclusion
add_bullet_slide("Conclusion", [
    "Senior/Expert roles dominate; Executive/Director rare",
    "Most jobs are full-time",
    "US leads in jobs and companies; Russia/US highest salaries",
    "Salaries rising with year and experience",
    "Medium/Large companies pay more",
    "Average salary: ~$112,298 USD"
])

# 14. Limitations
add_bullet_slide("Limitations", [
    "Insights based on 607 records",
    "May not represent the entire data science workforce"
])

# 15. Thank You
add_title_slide("Thank You", "Questions?")

prs.save("Data_Science_Job_Salaries_Presentation.pptx")
print("Presentation generated: Data_Science_Job_Salaries_Presentation.pptx")